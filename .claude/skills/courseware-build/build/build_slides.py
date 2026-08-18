#!/usr/bin/env python3
"""Generate the WSQ Application of BIM using Revit (TGS-2022015370) slide deck
(all-white Tertiary house style).

Design helpers are the same set used by the tertiary-course-slides skill that
produced the n8n reference deck (cover, section, content, two_col, cards3,
big_statement, step_slide, test_slide, brk). Content is driven entirely by
course_data.py + data_domainN.py + data_teaching.py (teaching slides and
graphics carried over from the legacy v9 master deck) so the deck stays 100%
aligned with the LP, LG and labs.
"""
import os, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, HERE)
import course_data as C
from data_domain1 import DOMAIN1
from data_domain2 import DOMAIN2
from data_domain3 import DOMAIN3
from data_domain4 import DOMAIN4
from data_domain5 import DOMAIN5
from data_teaching import TEACHING
ACTIVITIES = DOMAIN1 + DOMAIN2 + DOMAIN3 + DOMAIN4 + DOMAIN5

def _find_repo(start):
    """Locate the course repo (a dir containing both courseware/ and labs/).
    Env COURSE_REPO overrides. Keeps the build working wherever the skill lives."""
    env = os.environ.get("COURSE_REPO")
    if env and os.path.isdir(env):
        return env
    d = start
    for _ in range(8):
        d = os.path.dirname(d)
        if os.path.isdir(os.path.join(d, "courseware")) and os.path.isdir(os.path.join(d, "labs")):
            return d
    return os.path.dirname(os.path.dirname(HERE))
REPO = _find_repo(HERE)
ASSETS = os.path.join(os.path.dirname(HERE), "assets")   # co-located with the skill

# ---------------- palette (matches reference) ----------------
BLUE=RGBColor(0x1F,0x6F,0xEB); TEAL=RGBColor(0x10,0xB9,0x81); AMBER=RGBColor(0xF5,0x9E,0x0B)
INK=RGBColor(0x16,0x1B,0x26); GREY=RGBColor(0x5B,0x63,0x72); LIGHT=RGBColor(0xF5,0xF8,0xFC)
WHITE=RGBColor(0xFF,0xFF,0xFF); LINE=RGBColor(0xE2,0xE8,0xF0); VIOLET=RGBColor(0x7C,0x3A,0xED)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
SW,SH=prs.slide_width,prs.slide_height
BLANK=prs.slide_layouts[6]

def slide(): return prs.slides.add_slide(BLANK)
def rect(s,x,y,w,h,color,line=None):
    sp=s.shapes.add_shape(1,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(1)
    sp.shadow.inherit=False; return sp
def oval(s,x,y,w,h,color):
    sp=s.shapes.add_shape(9,x,y,w,h); sp.fill.solid(); sp.fill.fore_color.rgb=color
    sp.line.fill.background(); sp.shadow.inherit=False; return sp
def txt(s,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,space=4):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True; tf.vertical_anchor=anchor
    for i,line in enumerate(runs):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment=align; p.space_after=Pt(space)
        for t,sz,col,bold in line:
            r=p.add_run(); r.text=t; r.font.size=Pt(sz); r.font.bold=bold
            r.font.color.rgb=col; r.font.name="Arial"
    return tb
def bullets(s,x,y,w,h,items,size=18,color=INK,gap=10,mcolor=BLUE):
    tb=s.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    for i,it in enumerate(items):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.space_after=Pt(gap)
        lvl=it[1] if isinstance(it,tuple) else 0
        text=it[0] if isinstance(it,tuple) else it
        r=p.add_run(); r.text=("•  " if lvl==0 else "–  ")+text
        r.font.size=Pt(size if lvl==0 else size-2); r.font.color.rgb=color if lvl==0 else GREY
        r.font.name="Arial"; r.font.bold=(lvl==0 and isinstance(it,tuple) and len(it)>2 and it[2])
    return tb

PAGE={"n":0}
def footer(s):
    PAGE["n"]+=1
    txt(s,Inches(0.4),Inches(7.05),Inches(7.5),Inches(0.35),
        [[(f"{C.SHORT_TITLE}  ·  {C.COURSE_CODE}",9,GREY,False)]])
    txt(s,Inches(5.0),Inches(7.05),Inches(3.3),Inches(0.35),
        [[("© 2026 Tertiary Infotech Academy Pte Ltd",9,GREY,False)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(12.4),Inches(7.05),Inches(0.6),Inches(0.35),
        [[(str(PAGE["n"]),9,GREY,False)]],align=PP_ALIGN.RIGHT)
def head(s,title,kicker=None,kcolor=BLUE):
    rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),Inches(1.55),kcolor)
    if kicker: txt(s,Inches(0.85),Inches(0.5),Inches(11.6),Inches(0.4),[[(kicker,14,kcolor,True)]])
    txt(s,Inches(0.85),Inches(0.9),Inches(11.9),Inches(0.9),[[(title,29,INK,True)]])
    rect(s,Inches(0.85),Inches(1.7),Inches(11.63),Inches(0.02),LINE)
    return s
def _logo(name):
    p=os.path.join(ASSETS,name)
    return p if os.path.exists(p) else None

# ---------------- slide templates ----------------
def cover():
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),BLUE); rect(s,0,Inches(7.28),SW,Inches(0.22),TEAL)
    org=_logo("tertiary-infotech-logo.png")
    if org: s.shapes.add_picture(org,Inches(0.85),Inches(0.7),height=Inches(1.05))
    # course badge (top-right) — official CompTIA PenTest+ Certified badge, else text fallback
    badge=_logo("comptia-pentest-badge.png")
    if badge:
        s.shapes.add_picture(badge,Inches(10.05),Inches(0.6),width=Inches(2.5))
    else:
        rect(s,Inches(11.0),Inches(0.72),Inches(1.55),Inches(1.0),BLUE)
        txt(s,Inches(11.0),Inches(0.82),Inches(1.55),Inches(0.5),[[("PT0-003",18,WHITE,True)]],align=PP_ALIGN.CENTER)
        txt(s,Inches(11.0),Inches(1.28),Inches(1.55),Inches(0.4),[[("COMPTIA PENTEST+",8,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.9),Inches(2.3),Inches(12),Inches(0.6),[[("COURSE SLIDES  ·  WSQ",16,BLUE,True)]])
    txt(s,Inches(0.9),Inches(2.85),Inches(12.0),Inches(1.9),[[(C.TITLE,40,INK,True)]])
    rect(s,Inches(0.92),Inches(4.75),Inches(2.4),Inches(0.06),TEAL)
    txt(s,Inches(0.9),Inches(5.05),Inches(12),Inches(1.4),
        [[(f"WSQ Course Code: {C.COURSE_CODE}",16,GREY,False)],
         [("Conducted by Tertiary Infotech Academy Pte Ltd  ·  UEN 201200696W",14,GREY,False)]],space=6)
    txt(s,Inches(0.9),Inches(6.5),Inches(12),Inches(0.4),[[(f"Version {C.VERSION}  ·  {C.VERSION_DATE}",12,GREY,False)]])
    txt(s,Inches(0.9),Inches(6.85),Inches(12),Inches(0.34),[[("© 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved.  ·  www.tertiarycourses.com.sg",10,GREY,False)]])

def section(kicker,title,n,sub=""):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,BLUE)
    rect(s,Inches(0.85),Inches(2.5),Inches(0.14),Inches(2.0),TEAL)
    txt(s,Inches(1.25),Inches(2.55),Inches(11),Inches(0.6),[[(kicker,18,BLUE,True)]])
    txt(s,Inches(1.25),Inches(3.0),Inches(11.4),Inches(1.6),[[(title,40,INK,True)]])
    if sub: txt(s,Inches(1.27),Inches(4.55),Inches(11),Inches(0.8),[[(sub,16,GREY,False)]])
    txt(s,Inches(10.0),Inches(0.7),Inches(2.8),Inches(1.6),[[(n,72,RGBColor(0xE2,0xE8,0xF0),True)]],align=PP_ALIGN.RIGHT)
    footer(s)
def content(title,items,kicker=None,size=20):
    s=head(slide(),title,kicker); bullets(s,Inches(0.85),Inches(1.95),Inches(11.6),Inches(4.9),items,size=size); footer(s); return s
def two_col(title,left,right,kicker=None,lhead="",rhead=""):
    s=head(slide(),title,kicker)
    rect(s,Inches(0.85),Inches(1.95),Inches(5.7),Inches(4.7),LIGHT); rect(s,Inches(6.95),Inches(1.95),Inches(5.55),Inches(4.7),LIGHT)
    if lhead: txt(s,Inches(1.1),Inches(2.15),Inches(5.2),Inches(0.4),[[(lhead,16,BLUE,True)]])
    if rhead: txt(s,Inches(7.2),Inches(2.15),Inches(5.0),Inches(0.4),[[(rhead,16,TEAL,True)]])
    bullets(s,Inches(1.1),Inches(2.7),Inches(5.2),Inches(3.8),left,size=16)
    bullets(s,Inches(7.2),Inches(2.7),Inches(5.05),Inches(3.8),right,size=16,mcolor=TEAL); footer(s); return s
def cards3(title,cards,kicker):
    s=head(slide(),title,kicker); xs=[Inches(0.85),Inches(5.0),Inches(9.15)]
    for i,c in enumerate(cards[:3]):
        x=xs[i]; col=c[0]
        rect(s,x,Inches(1.95),Inches(3.65),Inches(4.7),LIGHT); rect(s,x,Inches(1.95),Inches(3.65),Inches(0.12),col)
        txt(s,x+Inches(0.25),Inches(2.2),Inches(3.2),Inches(0.6),[[(c[1],19,col,True)]])
        bullets(s,x+Inches(0.25),Inches(2.95),Inches(3.2),Inches(3.4),c[2],size=14,mcolor=col,gap=9)
    footer(s); return s
def big_statement(line1,line2,kicker,color=BLUE):
    s=slide(); rect(s,0,0,SW,SH,WHITE); rect(s,0,0,Inches(0.28),SH,color)
    txt(s,Inches(1.1),Inches(2.2),Inches(11),Inches(0.5),[[(kicker,16,color,True)]])
    txt(s,Inches(1.1),Inches(2.8),Inches(11.3),Inches(2.4),[[(line1,38,INK,True)]])
    if line2: txt(s,Inches(1.12),Inches(4.9),Inches(11),Inches(1.2),[[(line2,20,GREY,False)]])
    footer(s); return s
import math
PALETTE=[BLUE,TEAL,VIOLET,AMBER]
def tile_grid(title,items,kicker=None,cols=2,size=15,icons=None,accent=BLUE):
    """Grid of light panels, each with a coloured icon/number badge + text.
    items: list of strings (or (title,caption) tuples). Much richer than a bullet list."""
    s=head(slide(),title,kicker,kcolor=accent)
    n=len(items); rows=math.ceil(n/cols)
    X0=Inches(0.85); Y0=Inches(1.95); TOTW=Inches(11.63); AREAH=Inches(4.78)
    gx=Inches(0.3); gy=Inches(0.26)
    cw=int((TOTW-gx*(cols-1))/cols); ch=int((AREAH-gy*(rows-1))/rows)
    bd=Inches(0.6)
    for i,it in enumerate(items):
        r=i//cols; c=i%cols
        x=int(X0+(cw+gx)*c); y=int(Y0+(ch+gy)*r); col=PALETTE[i%len(PALETTE)]
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,Inches(0.1),ch,col)
        oval(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,col)
        ic=icons[i] if icons else str(i+1)
        txt(s,x+Inches(0.28),int(y+ch/2-bd/2),bd,bd,[[(ic,19,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        tx=x+Inches(1.08); tw=cw-Inches(1.32)
        if isinstance(it,tuple):
            txt(s,tx,int(y+Inches(0.14)),tw,int(ch-Inches(0.2)),
                [[(it[0],size+2,INK,True)],[(it[1],size-2,GREY,False)]],anchor=MSO_ANCHOR.MIDDLE,space=3)
        else:
            txt(s,tx,int(y+Inches(0.1)),tw,int(ch-Inches(0.16)),[[(it,size,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def flow_h(title,steps,kicker=None,color=BLUE):
    """Horizontal numbered flow: coloured chips connected by chevrons."""
    s=head(slide(),title,kicker,kcolor=color)
    n=len(steps); X0=Inches(0.85); TOTW=Inches(11.63); gap=Inches(0.34)
    cw=int((TOTW-gap*(n-1))/n); y=Inches(2.55); ch=Inches(3.15); bd=Inches(0.82)
    for i,st in enumerate(steps):
        x=int(X0+(cw+gap)*i)
        rect(s,x,y,cw,ch,LIGHT); rect(s,x,y,cw,Inches(0.1),color)
        oval(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,color)
        txt(s,int(x+cw/2-bd/2),int(y+Inches(0.42)),bd,bd,[[(str(i+1),30,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,x+Inches(0.16),int(y+Inches(1.55)),cw-Inches(0.32),int(ch-Inches(1.7)),[[(st,14,INK,False)]],align=PP_ALIGN.CENTER)
        if i<n-1:
            txt(s,int(x+cw-Inches(0.04)),int(y+ch/2-Inches(0.3)),int(gap+Inches(0.08)),Inches(0.6),
                [[("▶",15,color,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def trainer_slide(kicker,name,role,rows,initials,accent=BLUE):
    """Profile-card layout: avatar badge + name/role panel on the left, labelled
    info tiles on the right. rows: list of (LABEL, value); blank value → fill-in line."""
    s=head(slide(),"About the Trainer",kicker,kcolor=accent)
    lx=Inches(0.85); lw=Inches(3.65)
    rect(s,lx,Inches(1.95),lw,Inches(4.7),LIGHT); rect(s,lx,Inches(1.95),lw,Inches(0.12),accent)
    bd=Inches(1.7); ax=int(lx+(lw-bd)/2)
    oval(s,ax,Inches(2.5),bd,bd,accent)
    txt(s,ax,Inches(2.5),bd,bd,[[(initials,44,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
    txt(s,lx+Inches(0.15),Inches(4.55),lw-Inches(0.3),Inches(0.6),[[(name,21,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,lx+Inches(0.15),Inches(5.2),lw-Inches(0.3),Inches(1.2),[[(role,13,GREY,False)]],align=PP_ALIGN.CENTER)
    rx=Inches(4.9); rw=Inches(7.6); ry=Inches(1.95); rh=Inches(4.7)
    n=len(rows); gy=Inches(0.2); th=int((rh-gy*(n-1))/n)
    for i,(label,val) in enumerate(rows):
        y=int(ry+(th+gy)*i); col=PALETTE[i%len(PALETTE)]
        rect(s,rx,y,rw,th,LIGHT); rect(s,rx,y,Inches(0.1),th,col)
        vruns=[(val,14,INK,False)] if val else [("____________________________________________",13,LINE,False)]
        txt(s,rx+Inches(0.32),y,rw-Inches(0.6),th,
            [[(label.upper(),11,col,True)],vruns],anchor=MSO_ANCHOR.MIDDLE,space=3)
    footer(s); return s
def activity_overview(tag,title,desc,build,services,kicker):
    s=head(slide(),title,kicker,kcolor=TEAL)
    rect(s,Inches(0.85),Inches(1.85),Inches(1.7),Inches(0.5),TEAL)
    txt(s,Inches(0.85),Inches(1.9),Inches(1.7),Inches(0.4),[[(tag,16,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.85),Inches(2.55),Inches(11.7),Inches(1.6),[[(desc,21,INK,False)]])
    rect(s,Inches(0.85),Inches(4.3),Inches(11.7),Inches(2.0),LIGHT)
    txt(s,Inches(1.1),Inches(4.5),Inches(11),Inches(0.4),[[("You'll build",14,BLUE,True)]])
    txt(s,Inches(1.1),Inches(4.9),Inches(11),Inches(0.6),[[(build,18,INK,True)]])
    txt(s,Inches(1.1),Inches(5.6),Inches(11.2),Inches(0.6),[[("Tools:  ",13,GREY,True),(services,13,GREY,False)]]); footer(s); return s
def step_slide(kicker,act_title,n,total,text,cmd=""):
    s=head(slide(),act_title,kicker,TEAL)
    oval(s,Inches(0.85),Inches(2.5),Inches(1.4),Inches(1.4),TEAL)
    txt(s,Inches(0.85),Inches(2.74),Inches(1.4),Inches(0.9),[[(str(n),38,WHITE,True)]],align=PP_ALIGN.CENTER)
    txt(s,Inches(0.95),Inches(1.95),Inches(11),Inches(0.4),[[(f"STEP {n} OF {total}",13,GREY,True)]])
    txt(s,Inches(2.55),Inches(2.4),Inches(10.1),Inches(1.3),[[(text,23,INK,False)]],anchor=MSO_ANCHOR.MIDDLE)
    if cmd:
        rect(s,Inches(2.55),Inches(4.15),Inches(10.1),Inches(0.95),RGBColor(0x0B,0x12,0x20))
        txt(s,Inches(2.8),Inches(4.28),Inches(9.7),Inches(0.7),[[(cmd,13,RGBColor(0x9C,0xDC,0xFE),False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def steps2_slide(act_title,kicker,pairs,start_n,total):
    """Two lab steps stacked on one slide: number badge + instruction + Revit UI path chip."""
    s=head(slide(),act_title,kicker,TEAL)
    ys=[Inches(1.95),Inches(4.35)]
    for row,(n,instr,cmd) in enumerate(pairs):
        y=ys[row]
        rect(s,Inches(0.85),y,Inches(11.63),Inches(2.2),LIGHT)
        bd=Inches(0.95)
        oval(s,Inches(1.1),int(y+Inches(0.3)),bd,bd,TEAL)
        txt(s,Inches(1.1),int(y+Inches(0.3)),bd,bd,[[(str(n),26,WHITE,True)]],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE)
        txt(s,Inches(1.12),int(y+Inches(1.35)),Inches(1.0),Inches(0.35),[[(f"of {total}",10,GREY,True)]],align=PP_ALIGN.CENTER)
        txt(s,Inches(2.35),int(y+Inches(0.18)),Inches(9.9),Inches(1.3),[[(instr,14,INK,False)]],anchor=MSO_ANCHOR.TOP)
        if cmd:
            rect(s,Inches(2.35),int(y+Inches(1.6)),Inches(9.9),Inches(0.48),RGBColor(0x0B,0x12,0x20))
            txt(s,Inches(2.55),int(y+Inches(1.63)),Inches(9.55),Inches(0.42),[[(cmd,11,RGBColor(0x9C,0xDC,0xFE),False)]],anchor=MSO_ANCHOR.MIDDLE)
    footer(s); return s
def test_slide(act_title,text,kicker):
    s=head(slide(),act_title,kicker,TEAL)
    rect(s,Inches(0.85),Inches(2.3),Inches(11.7),Inches(2.6),RGBColor(0xE8,0xF7,0xEE))
    txt(s,Inches(1.2),Inches(2.6),Inches(11),Inches(0.5),[[("✅  Test it",20,RGBColor(0x12,0x7A,0x3E),True)]])
    txt(s,Inches(1.2),Inches(3.3),Inches(11),Inches(1.4),[[(text,18,INK,False)]]); footer(s); return s
def brk(kind,dur,color=AMBER):
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),color); rect(s,0,Inches(7.28),SW,Inches(0.22),color)
    rect(s,Inches(5.4),Inches(2.35),Inches(2.53),Inches(0.1),color)
    txt(s,0,Inches(2.75),SW,Inches(1.2),[[(kind,48,INK,True)]],align=PP_ALIGN.CENTER)
    txt(s,0,Inches(4.05),SW,Inches(0.8),[[(dur,22,color,True)]],align=PP_ALIGN.CENTER); PAGE["n"]+=1

# ============================================================ BUILD
LEGACY=os.path.join(REPO,"courseware","assets","legacy")
def _legacy(name):
    p=os.path.join(LEGACY,name)
    return p if os.path.exists(p) else None
try:
    from PIL import Image
    def _imgsize(p):
        with Image.open(p) as im: return im.size
except Exception:
    def _imgsize(p): return (1000,600)

def _fit(p,maxw,maxh):
    """Return (w,h) EMU fitting the image inside maxw×maxh (inches)."""
    w0,h0=_imgsize(p); ar=h0/float(w0)
    w=maxw; h=w*ar
    if h>maxh: h=maxh; w=h/ar
    return Inches(w),Inches(h)

def pic_slide(title,kicker,items,image,imgw=0):
    """Bullets left · legacy image right (framed)."""
    s=head(slide(),title,kicker)
    p=_legacy(image)
    iw=imgw if imgw else 4.8
    if p:
        w,h=_fit(p,iw,4.55)
        x=int(Inches(12.48)-w); y=Inches(2.05)
        rect(s,int(x-Inches(0.08)),int(y-Inches(0.08)),int(w+Inches(0.16)),int(h+Inches(0.16)),LIGHT)
        s.shapes.add_picture(p,x,y,width=w,height=h)
        bw=Inches(12.48)-w-Inches(1.25)
    else:
        bw=Inches(11.6)
    bullets(s,Inches(0.85),Inches(2.0),int(bw),Inches(4.7),items,size=15,gap=12)
    footer(s); return s

def img_slide(title,kicker,image,caption=""):
    """Full-width legacy image with caption."""
    s=head(slide(),title,kicker)
    p=_legacy(image)
    if p:
        w,h=_fit(p,11.3,4.35)
        x=int((SW-w)/2); y=Inches(2.0)
        rect(s,int(x-Inches(0.08)),int(y-Inches(0.08)),int(w+Inches(0.16)),int(h+Inches(0.16)),LIGHT)
        s.shapes.add_picture(p,x,y,width=w,height=h)
    if caption:
        txt(s,Inches(0.85),Inches(6.5),Inches(11.6),Inches(0.5),[[(caption,13,GREY,False)]],align=PP_ALIGN.CENTER)
    footer(s); return s

def banner_slide(title,kicker,items,image):
    """Wide ribbon-strip image on top, bullets below."""
    s=head(slide(),title,kicker)
    p=_legacy(image)
    y=Inches(1.95)
    if p:
        w,h=_fit(p,11.63,2.05)
        x=int((SW-w)/2)
        rect(s,int(x-Inches(0.06)),int(y-Inches(0.06)),int(w+Inches(0.12)),int(h+Inches(0.12)),LIGHT)
        s.shapes.add_picture(p,x,y,width=w,height=h)
        by=int(y+h+Inches(0.3))
    else:
        by=Inches(2.1)
    bullets(s,Inches(0.85),by,Inches(11.6),int(Inches(6.9)-by),items,size=14,gap=10)
    footer(s); return s

def render_teaching(tnum,tcode):
    for item in TEACHING.get(tnum,[]):
        kind=item[0]
        if kind=="tiles":
            _,title,kicker,items,cols,size=item
            tile_grid(title,items,kicker=kicker,cols=cols,size=size)
        elif kind=="pic":
            _,title,kicker,items,image,imgw=item
            pic_slide(title,kicker,items,image,imgw)
        elif kind=="img":
            _,title,kicker,image,caption=item
            img_slide(title,kicker,image,caption)
        elif kind=="banner":
            _,title,kicker,items,image=item
            banner_slide(title,kicker,items,image)
        elif kind=="two":
            _,title,kicker,left,right,lhead,rhead=item
            two_col(title,left,right,kicker=kicker,lhead=lhead,rhead=rhead)
        elif kind=="flow":
            _,title,kicker,steps=item
            flow_h(title,steps,kicker=kicker)
        elif kind=="big":
            _,l1,l2,kicker=item
            big_statement(l1,l2,kicker)

# --- course cover (BIM hero image, no cert badge) ---
def bim_cover():
    s=slide(); rect(s,0,0,SW,SH,WHITE)
    rect(s,0,0,SW,Inches(0.22),BLUE); rect(s,0,Inches(7.28),SW,Inches(0.22),TEAL)
    org=_logo("tertiary-infotech-logo.png")
    if org: s.shapes.add_picture(org,Inches(0.85),Inches(0.7),height=Inches(1.05))
    hero=_legacy("slide001_img1.png")
    if hero:
        w,h=_fit(hero,3.6,3.4)
        x=int(Inches(12.45)-w)
        rect(s,int(x-Inches(0.08)),Inches(2.32),int(w+Inches(0.16)),int(h+Inches(0.16)),LIGHT)
        s.shapes.add_picture(hero,x,Inches(2.4),width=w,height=h)
    txt(s,Inches(0.9),Inches(2.3),Inches(8.5),Inches(0.6),[[("COURSE SLIDES  ·  WSQ",16,BLUE,True)]])
    txt(s,Inches(0.9),Inches(2.85),Inches(8.6),Inches(1.9),[[(C.TITLE,40,INK,True)]])
    rect(s,Inches(0.92),Inches(4.35),Inches(2.4),Inches(0.06),TEAL)
    txt(s,Inches(0.9),Inches(4.7),Inches(8.5),Inches(1.6),
        [[(f"WSQ Course Code: {C.COURSE_CODE}",16,GREY,False)],
         [(f"Skills Framework TSC: {C.TSC_TITLE} ({C.TSC_CODE})",13,GREY,False)],
         [("Conducted by Tertiary Infotech Academy Pte Ltd  ·  UEN 201200696W",14,GREY,False)]],space=6)
    txt(s,Inches(0.9),Inches(6.45),Inches(12),Inches(0.4),[[(f"Version {C.VERSION}  ·  {C.VERSION_DATE}",12,GREY,False)]])
    txt(s,Inches(0.9),Inches(6.85),Inches(12),Inches(0.34),[[("© 2026 Tertiary Infotech Academy Pte Ltd. All rights reserved.  ·  www.tertiarycourses.com.sg",10,GREY,False)]])
bim_cover()

# ---------------- ADMIN (front) ----------------
section("COURSE ADMINISTRATION","Welcome & Housekeeping","")
content("Digital Attendance (Mandatory)",[
 "It is mandatory to take the AM, PM and Assessment digital attendance for WSQ-funded courses.",
 "The trainer/administrator displays the digital attendance QR code from the SSG portal.",
 "Scan the QR code with your mobile phone camera and submit your attendance.",
 "A minimum of 75% attendance is required to be eligible for assessment and funding."],kicker="TRAQOM · SSG DIGITAL ATTENDANCE")
trainer_slide("YOUR TRAINER · GENERAL","Your Trainer","General Trainer template —\nto be completed by the trainer",
 [("Name",""),("Title / Designation",""),("Qualifications",""),
  ("Areas of expertise",""),("Training & industry experience",""),("Contact","")],
 initials="?",accent=GREY)
trainer_slide("YOUR TRAINER",C.TRAINER,"Principal Trainer\nTertiary Infotech Academy Pte Ltd",
 [("Role","Principal Trainer, Tertiary Infotech Academy Pte Ltd"),
  ("Expertise","BIM and parametric design — Autodesk Revit, AutoCAD, Fusion 360 and engineering CAD/CAE toolchains."),
  ("Delivers","WSQ courses on BIM with Revit, architecture drawing, technical drawing and product design."),
  ("Founder","Founder and lead instructor at Tertiary Infotech / Tertiary Courses.")],
 initials="AA",accent=BLUE)
content("Let's Know Each Other",[
 "Your name and organisation / role.",
 "Your experience with CAD, Revit or the built-environment industry (if any).",
 "What you want to be able to model, manage or submit with BIM after this course."],kicker="ICE-BREAKER")
tile_grid("Ground Rules",[
 "Set your mobile phone to silent mode.","Participate actively — no question is too small.",
 "Mutual respect: agree to disagree.","One conversation at a time.",
 "Be punctual; return from breaks on time.","75% attendance is required."],
 kicker="HOUSEKEEPING",cols=2,size=15)
pic_slide("Download Course Material","LMS / TMS · lms-tms.tertiaryinfotech.com",[
 "Log in to the LMS/TMS portal at https://lms-tms.tertiaryinfotech.com with the account created for your enrolment.",
 "Open your course — WSQ Application of BIM using Revit — from My Courses.",
 "Download the course slides, Learner Guide and the Revit lab dataset files (.rvt) before the hands-on sessions.",
 "The same portal carries your digital attendance records, assessment submission and course feedback (TRAQOM).",
 ],"slide015_img3.png",4.0)
two_col("Lesson Plan — 2 Days, 8 Hours/Day",[
 (f"Day 1 — {C.DAY_THEMES[1]}",0),
 ("Digital attendance (AM) · welcome & introductions",1),
 ("Topic 1: Introduction to BIM and Revit (Labs 1–2)",1),
 ("Lunch · digital attendance (PM)",1),
 ("Topic 2: BIM Modeling (Labs 3–6)",1),
 ("Day 1 recap and Q&A",1)],
 [(f"Day 2 — {C.DAY_THEMES[2]}",0),
 ("Digital attendance (AM) · Day 1 recap",1),
 ("Topic 3: BIM Application (Labs 7–8)",1),
 ("Lunch · digital attendance (PM)",1),
 ("Topic 4: BIM Documentation (Labs 9–10)",1),
 ("Final Assessment: WA (SAQ) + Practical Performance",1),
 ("Daily timing: 9:30am–6:30pm · 1-hour lunch · tea breaks within",1)],
 kicker="SCHEDULE",lhead="Day 1",rhead="Day 2")
tile_grid("Learning Outcomes",[
 ("LO1","Apply Revit for Building Information Modelling (BIM)."),
 ("LO2","Develop models and integrate design for BIM."),
 ("LO3","Apply BIM and check the output."),
 ("LO4","Maintain BIM database and develop BIM documentation.")],
 kicker="WHAT YOU'LL ACHIEVE",cols=2,size=15)
big_statement("Skills Framework for the Built Environment",
 f"TSC: {C.TSC_TITLE}   ·   TSC Code: {C.TSC_CODE}","SKILLS FRAMEWORK ALIGNMENT")
tile_grid("TSC Knowledge — What the WA (SAQ) Assesses",
 [(k,v) for k,v in C.TSC_KNOWLEDGE],
 kicker=f"{C.TSC_CODE} · KNOWLEDGE K1–K10",cols=2,size=12)
tile_grid("TSC Abilities — What the Practical Performance Assesses",
 [(a,v) for a,v in C.TSC_ABILITIES],
 kicker=f"{C.TSC_CODE} · ABILITIES A1–A9",cols=2,size=12)
tile_grid("Course Outline — Four Topics, Ten Hands-On Labs",[
 ("Topic 1 · Introduction to BIM and Revit","K1–K4, A1–A2 — BIM principles, dimensions, value, requirements; the Revit environment. Labs 1–2."),
 ("Topic 2 · BIM Modeling","K6–K7, A4–A5 — walls, floors, roofs, stairs, structure, site, massing, families. Labs 3–6."),
 ("Topic 3 · BIM Application","K5, A6–A8 — operate BIM, schedules and takeoffs, energy analysis, compliance checks. Labs 7–8."),
 ("Topic 4 · BIM Documentation","K8–K10, A3, A9 — documentation, databases, worksharing, BCA e-submission standards. Labs 9–10.")],
 kicker="COURSE MAP",cols=2,size=13)
pic_slide("Practice Exam","TEST YOURSELF · exams.tertiaryinfotech.com",[
 "Sharpen your readiness with the Tertiary Infotech practice exams portal.",
 "Portal: https://exams.tertiaryinfotech.com — search for the BIM / Revit practice exam.",
 "Attempt it under timed conditions and review every explanation.",
 "Revisit any topic you miss, then re-take the practice exam before the final assessment.",
 ],"slide184_img96.png",4.6)
content("Briefing for Assessment",[
 "Place phones and other materials under the table or on the floor.",
 "No photos or recording of assessment scripts.","No discussion during the assessment.",
 "Use a black/blue pen for hard-copy assessments.","No liquid paper / correction tape.",
 "Scripts are collected when time is up."])
content("Assessment",[
 C.ASSESSMENT["written"], C.ASSESSMENT["practical"],
 "Format: Open Book — slides, Learner Guide and approved materials only.",
 "The PP is completed in Revit using the provided PP dataset (.rvt).",
 C.ASSESSMENT["note"],"An appeal process is available if required."],kicker="FINAL ASSESSMENT · DAY 2")
flow_h("Assessment Flow",[
 "TRAQOM survey — scan the QR code on the LMS",
 "Assessment digital attendance — scan the SSG QR",
 "Sit WA (SAQ, 1 hr) then PP (1.5 hrs) — open book",
 "Submit your answers on the LMS",
 "Sign the Assessment Summary Record"],kicker="ON ASSESSMENT DAY")

# ---------------- TOPICS: teaching + labs ----------------
TOPIC_ACTS = {t["num"]: [a for a in ACTIVITIES if a["topic"]==t["num"]] for t in C.TOPICS}
CARD_COLORS=[BLUE,TEAL,VIOLET]
for t in C.TOPICS:
    section(f"TOPIC {t['code']}", t["title"], t["code"], t["subtitle"])
    tile_grid(f"Key Concepts — {t['title']}", t["concepts"],
              kicker=f"ALIGNMENT: {t['weighting']}", cols=2, size=13)
    render_teaching(t["num"],t["code"])
    acts=TOPIC_ACTS[t["num"]]
    if len(acts)>1:
        half=(len(acts)+1)//2
        groups=[acts[:half],acts[half:],[]]
        cards=[]
        for gi,g in enumerate(groups[:3]):
            if g: cards.append((CARD_COLORS[gi], f"Lab {g[0]['num']}" if len(g)==1 else f"Labs {g[0]['num']}–{g[-1]['num']}", [a["title"] for a in g]))
        while len(cards)<3: cards.append((CARD_COLORS[len(cards)],"—",["—"]))
        cards3(f"Hands-On Labs — {t['title']}", cards, kicker="WHAT YOU'LL DO")
    for a in acts:
        activity_overview(f"LAB {a['num']}", a["title"], a["desc"], a["build"], a["services"], kicker=f"TOPIC {t['code']} · HANDS-ON")
        steps=a["steps"]; total=len(steps)
        for i in range(0,total,2):
            pairs=[(j+1,steps[j][0],steps[j][1]) for j in range(i,min(i+2,total))]
            steps2_slide(a["title"],f"LAB {a['num']} · STEP-BY-STEP",pairs,i+1,total)
        test_slide(a["title"], a["test"], kicker=f"LAB {a['num']} · VERIFY")
    content(f"Recap — {t['title']}",
            ["You can now: "+o for o in dict.fromkeys(a["objective"] for a in acts)][:6],
            kicker="TOPIC RECAP", size=17)

# ---------------- CLOSE ----------------
section("WRAP-UP","Course Summary & Next Steps","")
tile_grid("What You Achieved",[
 ("Applied Revit for BIM","Navigated the Revit BIM environment and turned 2D information into an intelligent 3D model (LO1)."),
 ("Developed & integrated models","Modelled walls, floors, roofs, stairs, structure, site and masses — with families and materials (LO2)."),
 ("Applied BIM & checked output","Built schedules and takeoffs, ran energy analysis and interference checks on the model (LO3)."),
 ("Maintained database & documentation","Workshared a central model, produced sheets, and met BCA e-submission standards (LO4).")],
 kicker="LEARNING OUTCOMES",cols=2,size=14)
pic_slide("Certificate & TRAQOM Survey (Mandatory)","AFTER THE COURSE · lms-tms.tertiaryinfotech.com",[
 "Complete the TRAQOM course-quality survey on the LMS — it is mandatory for WSQ funding.",
 "Your WSQ Statement of Attainment (SOA) is issued after you are assessed as Competent.",
 "Download your certificate of completion from the LMS portal.",
 ],"slide184_img96.png",5.2)
tile_grid("Recommended Next Courses",[(rc,"") for rc in C.RECOMMENDED_COURSES]+[("Browse all courses","www.tertiarycourses.com.sg")],
 kicker="KEEP LEARNING",cols=2,size=13)
content("Support",[
 "If you have any enquiries during and after the class, contact us:",
 "Email: enquiry@tertiaryinfotech.com",
 "Tel: +65 6100 0613",
 "Website: www.tertiarycourses.com.sg"],kicker="WE'RE HERE TO HELP")
content("Assessment",[
 C.ASSESSMENT["written"], C.ASSESSMENT["practical"],
 "Open book: slides, Learner Guide and approved materials only.",
 "Remember to take the Assessment digital attendance (TRAQOM).",
 "Submit your completed answers on the LMS at https://lms-tms.tertiaryinfotech.com/."],kicker="FINAL ASSESSMENT")
flow_h("Assessment Flow",[
 "TRAQOM survey — scan the QR code on the LMS",
 "Assessment digital attendance — scan the SSG QR",
 "Sit WA (SAQ, 1 hr) then PP (1.5 hrs) — open book",
 "Submit your answers on the LMS",
 "Sign the Assessment Summary Record"],kicker="ON ASSESSMENT DAY")
content("Digital Attendance (Mandatory)",[
 "It is mandatory to take the AM, PM and Assessment digital attendance for WSQ-funded courses.",
 "The trainer/administrator displays the digital attendance QR code from the SSG portal.",
 "Scan the QR code with your mobile phone camera and submit your attendance.",
 "A minimum of 75% attendance is required to be eligible for assessment and funding."],kicker="TRAQOM · SSG DIGITAL ATTENDANCE")
big_statement("Thank You!","You are now ready to apply, model, manage and document BIM projects with Autodesk Revit.","HAPPY MODELLING",color=TEAL)

OUT=os.path.join(REPO,"courseware",f"{C.SHORT_TITLE}-{C.VERSION}.pptx")
prs.save(OUT)
print(f"Saved {OUT}  ({len(prs.slides._sldIdLst)} slides)")
